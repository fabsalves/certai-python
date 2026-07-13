"""Plain-text extraction from uploaded documents (txt/docx/pdf/pptx)."""

import io

# Truncation cap before handing the text to the LLM.
MAX_EXTRACTED_CHARS = 60_000


class UnsupportedFormatError(Exception):
    """Format without a reliable pure-Python extractor (e.g. legacy .ppt)."""


def extract_text(content: bytes, extension: str) -> str:
    ext = (extension or "").lower()
    if ext == ".txt":
        text = content.decode("utf-8", errors="replace")
    elif ext == ".docx":
        text = _extract_docx(content)
    elif ext == ".pdf":
        text = _extract_pdf(content)
    elif ext == ".pptx":
        text = _extract_pptx(content)
    else:
        raise UnsupportedFormatError(
            f"Sem extrator de texto para {ext or 'arquivo sem extensão'}"
        )
    return text.strip()[:MAX_EXTRACTED_CHARS]


def _extract_docx(content: bytes) -> str:
    from docx import Document

    document = Document(io.BytesIO(content))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"Notas: {notes}")
        if slide_texts:
            parts.append(f"[Slide {index}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)
