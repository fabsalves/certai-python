#!/usr/bin/env python3
"""Gera os materiais de demonstração em docs/materiais-demo/."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Pt

OUT_DIR = Path(__file__).resolve().parent

TRACK_SLIDES: list[tuple[str, list[str]]] = [
    (
        "Comunicação escrita no trabalho",
        [
            "Programa da trilha — material de referência para professores e tutores.",
            "Competência macro: redigir pareceres claros e objetivos no contexto profissional.",
        ],
    ),
    (
        "Sequência — Fundamentos",
        [
            "Leitura crítica de textos",
            "Estrutura de um parecer",
            "Primeiro rascunho",
        ],
    ),
    (
        "Sequência — Prática",
        [
            "Revisão em pares",
            "Argumentação objetiva",
            "Entrega final",
        ],
    ),
    (
        "Papel de cada ator",
        [
            "Designer: publica o conteúdo fixo de cada aula na trilha.",
            "Professor: conduz a aula em sala e encerra com relato + material usado.",
            "Lira: conversa com cada aluno sobre o que a turma já estudou e extrai evidência de entendimento.",
        ],
    ),
    (
        "Critérios transversais",
        [
            "Separar fato de interpretação antes de opinar sobre intenção.",
            "Parecer em três blocos: contexto, análise, recomendação.",
            "Rascunho primeiro; revisão focada em clareza para o leitor.",
            "Argumentação amarrada a evidência do caso.",
        ],
    ),
    (
        "Aula 1 em foco",
        [
            "Objetivo: leitura crítica de e-mails e mensagens de trabalho.",
            'Pergunta-guia: "Isso está escrito ou estou supondo?"',
            'Texto central: "Não vou aprovar isso agora."',
            "A Lira conduz aplicação do material liberado — não ensina aulas futuras.",
        ],
    ),
]


def _heading(doc: Document, text: str) -> None:
    doc.add_heading(text, level=1)


def _subheading(doc: Document, text: str) -> None:
    doc.add_heading(text, level=2)


def _body(doc: Document, text: str) -> None:
    for paragraph in text.strip().split("\n"):
        doc.add_paragraph(paragraph)


def _add_bullets(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(18 if index == 0 and len(lines) == 1 else 16)


def build_track_pptx() -> Path:
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    for index, (title, bullets) in enumerate(TRACK_SLIDES):
        layout = title_layout if index == 0 else bullet_layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if index == 0:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = bullets[0]
        else:
            body = slide.placeholders[1].text_frame
            _add_bullets(body, bullets)

    path = OUT_DIR / "trilha-comunicacao-escrita-trabalho.pptx"
    prs.save(path)
    return path


def _dejavu_font_path() -> Path:
    candidates = [
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/TTF/DejaVuSans.ttf"),
    ]
    for path in candidates:
        if path.is_file():
            return path
    raise SystemExit(
        "DejaVuSans.ttf não encontrado no sistema. "
        "Use o PPTX da trilha ou instale fonts-dejavu-core."
    )


def build_track_pdf() -> Path:
    try:
        from fpdf import FPDF
    except ImportError as exc:
        raise SystemExit(
            "fpdf2 é necessário para gerar o PDF da trilha. "
            "Instale com: backend/.venv/bin/pip install fpdf2"
        ) from exc

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_margins(15, 15, 15)
    pdf.add_page()
    content_width = pdf.epw
    font_regular = _dejavu_font_path()
    font_bold = font_regular.with_name("DejaVuSans-Bold.ttf")
    pdf.add_font("DejaVu", "", str(font_regular))
    if font_bold.is_file():
        pdf.add_font("DejaVu", "B", str(font_bold))
    pdf.set_font("DejaVu", size=11)

    for title, bullets in TRACK_SLIDES:
        style = "B" if font_bold.is_file() else ""
        pdf.set_font("DejaVu", style, 14)
        pdf.multi_cell(content_width, 8, title)
        pdf.ln(2)
        pdf.set_font("DejaVu", size=11)
        for line in bullets:
            pdf.multi_cell(content_width, 6, f"- {line}")
        pdf.ln(4)

    path = OUT_DIR / "trilha-comunicacao-escrita-trabalho.pdf"
    pdf.output(path)
    return path


def build_lesson1_attachment() -> Path:
    doc = Document()
    _heading(doc, "Aula 1 — Leitura crítica de textos")
    _body(doc, "Material usado em sala (complemento ao conteúdo publicado na trilha).")

    _subheading(doc, "Definições")
    _body(
        doc,
        """
Fato: o que está literalmente no texto — dá para apontar a frase.
Interpretação: o que você conclui a partir do que leu.
Pergunta-guia: "Isso está escrito ou estou supondo?"
        """,
    )

    _subheading(doc, "Exemplo trabalhado em sala")
    _body(
        doc,
        """
E-mail interno:
"Não vou aprovar isso agora."

Atividade: em duplas, listem o que é fato e o que seria interpretação.
Atenção: "não aprovar agora" não é o mesmo que "nunca vai aprovar".
Tom irritado ou cordial não está escrito — é suposição se não houver evidência no texto.
        """,
    )

    _subheading(doc, "Segundo exemplo (prática em sala)")
    _body(
        doc,
        """
Mensagem de equipe:
"O relatório precisa ser revisado até amanhã. O cliente parece insatisfeito com os resultados."

Perguntas para a turma:
- O que está escrito literalmente?
- O que é conclusão sobre o cliente ou sobre urgência?
        """,
    )

    _subheading(doc, "Confusões que apareceram na turma")
    _body(
        doc,
        """
- Tratar tom da mensagem (secos, frios, educados) como se fosse fato.
- Transformar "não agora" em "nunca" ou em "decisão final".
- Misturar recomendação pessoal ("eu mandaria outro e-mail") com leitura do texto.
        """,
    )

    _subheading(doc, "Para a conversa com a Lira")
    _body(
        doc,
        """
Cada aluno deve:
1. Classificar o e-mail "Não vou aprovar isso agora." com as próprias palavras.
2. Trazer um trecho curto de mensagem ou e-mail do trabalho e repetir o exercício.
        """,
    )

    path = OUT_DIR / "aula-01-leitura-critica-anexo.docx"
    doc.save(path)
    return path


def main() -> None:
    paths = [build_track_pptx(), build_track_pdf(), build_lesson1_attachment()]
    for path in paths:
        print(f"Wrote {path}")


if __name__ == "__main__":
    main()
